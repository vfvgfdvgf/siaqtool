from __future__ import annotations

import csv
from pathlib import Path

import pdfplumber
import pymupdf as fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt

from .core import ProcessResult, ProcessingError, file_result, libreoffice_convert, pdf_result


FORMAT_TARGETS = {
    "doc-to-docx": ("docx", "siaq-document.docx"),
    "docx-to-doc": ("doc", "siaq-document.doc"),
    "xls-to-xlsx": ("xlsx", "siaq-workbook.xlsx"),
    "xlsx-to-xls": ("xls", "siaq-workbook.xls"),
    "ppt-to-pptx": ("pptx", "siaq-presentation.pptx"),
    "pptx-to-ppt": ("ppt", "siaq-presentation.ppt"),
    "odt-to-docx": ("docx", "siaq-document.docx"),
    "docx-to-odt": ("odt", "siaq-document.odt"),
    "ods-to-xlsx": ("xlsx", "siaq-workbook.xlsx"),
    "xlsx-to-ods": ("ods", "siaq-workbook.ods"),
    "odp-to-pptx": ("pptx", "siaq-presentation.pptx"),
    "pptx-to-odp": ("odp", "siaq-presentation.odp"),
    "rtf-to-docx": ("docx", "siaq-document.docx"),
    "docx-to-rtf": ("rtf", "siaq-document.rtf"),
}


def office_to_pdf(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    result = libreoffice_convert(paths[0], workdir, "pdf")
    return pdf_result(result, "siaq-converted.pdf")


def office_format(paths: list[Path], workdir: Path, _options: dict, slug: str) -> ProcessResult:
    target, filename = FORMAT_TARGETS[slug]
    result = libreoffice_convert(paths[0], workdir, target)
    return file_result(result, filename)


def _modern_word(path: Path, workdir: Path) -> Path:
    return path if path.suffix.lower() == ".docx" else libreoffice_convert(path, workdir, "docx")


def _modern_excel(path: Path, workdir: Path) -> Path:
    if path.suffix.lower() == ".xlsx":
        return path
    if path.suffix.lower() == ".csv":
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "البيانات"
        text = path.read_text(encoding="utf-8-sig", errors="replace")
        for row in csv.reader(text.splitlines()):
            sheet.append(row)
        output = workdir / "source.xlsx"
        workbook.save(output)
        return output
    return libreoffice_convert(path, workdir, "xlsx")


def _modern_powerpoint(path: Path, workdir: Path) -> Path:
    return path if path.suffix.lower() == ".pptx" else libreoffice_convert(path, workdir, "pptx")


def pdf_to_word(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    source = fitz.open(paths[0])
    document = Document()
    for index, page in enumerate(source, start=1):
        if index > 1:
            document.add_page_break()
        document.add_heading(f"الصفحة {index}", level=1)
        blocks = page.get_text("blocks")
        for block in sorted(blocks, key=lambda item: (round(item[1]), item[0])):
            text = str(block[4]).strip()
            if text:
                document.add_paragraph(text)
    source.close()
    output = workdir / "document.docx"
    document.save(output)
    return file_result(output, "siaq-pdf-to-word.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


def pdf_to_excel(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    workbook = Workbook()
    workbook.remove(workbook.active)
    with pdfplumber.open(paths[0]) as document:
        for page_index, page in enumerate(document.pages, start=1):
            sheet = workbook.create_sheet(f"صفحة {page_index}")
            tables = page.extract_tables() or []
            row_cursor = 1
            for table in tables:
                for row in table:
                    for column_index, value in enumerate(row or [], start=1):
                        sheet.cell(row_cursor, column_index, value or "")
                    row_cursor += 1
                row_cursor += 1
            if not tables:
                for row_index, line in enumerate((page.extract_text() or "").splitlines(), start=1):
                    sheet.cell(row_index, 1, line)
    if not workbook.sheetnames:
        workbook.create_sheet("النتيجة")
    output = workdir / "document.xlsx"
    workbook.save(output)
    return file_result(output, "siaq-pdf-to-excel.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def pdf_to_powerpoint(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    document = fitz.open(paths[0])
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for index, page in enumerate(document, start=1):
        image_path = workdir / f"slide-{index}.png"
        page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).save(image_path)
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    document.close()
    output = workdir / "document.pptx"
    presentation.save(output)
    return file_result(output, "siaq-pdf-to-powerpoint.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")


def excel_to_word(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    workbook = load_workbook(_modern_excel(paths[0], workdir), data_only=True, read_only=True)
    document = Document()
    for sheet_index, sheet in enumerate(workbook.worksheets):
        if sheet_index:
            document.add_page_break()
        document.add_heading(sheet.title, level=1)
        rows = list(sheet.iter_rows(values_only=True, max_row=min(sheet.max_row, 300), max_col=min(sheet.max_column, 40)))
        width = max((len(row) for row in rows), default=1)
        table = document.add_table(rows=max(len(rows), 1), cols=max(width, 1))
        table.style = "Table Grid"
        for row_index, row in enumerate(rows):
            for column_index, value in enumerate(row):
                table.cell(row_index, column_index).text = "" if value is None else str(value)
    output = workdir / "excel.docx"
    document.save(output)
    return file_result(output, "siaq-excel-to-word.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


def word_to_excel(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    document = Document(_modern_word(paths[0], workdir))
    workbook = Workbook()
    paragraphs = workbook.active
    paragraphs.title = "النص"
    for row, paragraph in enumerate((item.text for item in document.paragraphs if item.text.strip()), start=1):
        paragraphs.cell(row, 1, paragraph)
    for table_index, table in enumerate(document.tables, start=1):
        sheet = workbook.create_sheet(f"جدول {table_index}")
        for row_index, row in enumerate(table.rows, start=1):
            for column_index, cell in enumerate(row.cells, start=1):
                sheet.cell(row_index, column_index, cell.text)
    output = workdir / "word.xlsx"
    workbook.save(output)
    return file_result(output, "siaq-word-to-excel.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def word_to_powerpoint(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    document = Document(_modern_word(paths[0], workdir))
    presentation = Presentation()
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    chunks = [paragraphs[index:index + 6] for index in range(0, len(paragraphs), 6)] or [["مستند فارغ"]]
    for index, chunk in enumerate(chunks, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = chunk[0][:120] if chunk else f"الشريحة {index}"
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for item in chunk[1:]:
            paragraph = frame.add_paragraph()
            paragraph.text = item[:500]
            paragraph.font.size = Pt(22)
    output = workdir / "word.pptx"
    presentation.save(output)
    return file_result(output, "siaq-word-to-powerpoint.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")


def powerpoint_to_word(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    presentation = Presentation(_modern_powerpoint(paths[0], workdir))
    document = Document()
    for index, slide in enumerate(presentation.slides, start=1):
        document.add_heading(f"الشريحة {index}", level=1)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                document.add_paragraph(shape.text.strip())
    output = workdir / "powerpoint.docx"
    document.save(output)
    return file_result(output, "siaq-powerpoint-to-word.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


def excel_to_powerpoint(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    workbook = load_workbook(_modern_excel(paths[0], workdir), data_only=True, read_only=True)
    presentation = Presentation()
    for sheet in workbook.worksheets:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = sheet.title
        lines = []
        for row in sheet.iter_rows(values_only=True, max_row=min(sheet.max_row, 20), max_col=min(sheet.max_column, 8)):
            lines.append(" | ".join("" if value is None else str(value) for value in row))
        slide.placeholders[1].text = "\n".join(lines)[:5000] or "لا توجد بيانات"
    output = workdir / "excel.pptx"
    presentation.save(output)
    return file_result(output, "siaq-excel-to-powerpoint.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")


def powerpoint_to_excel(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    presentation = Presentation(_modern_powerpoint(paths[0], workdir))
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "الشرائح"
    sheet.append(["الشريحة", "النص"])
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        sheet.append([index, "\n".join(texts)])
    output = workdir / "powerpoint.xlsx"
    workbook.save(output)
    return file_result(output, "siaq-powerpoint-to-excel.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def csv_to_excel(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    output = _modern_excel(paths[0], workdir)
    return file_result(output, "siaq-csv-to-excel.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def excel_to_csv(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    workbook = load_workbook(_modern_excel(paths[0], workdir), data_only=True, read_only=True)
    output = workdir / "sheet.csv"
    with output.open("w", encoding="utf-8-sig", newline="") as target:
        writer = csv.writer(target)
        for row in workbook.active.iter_rows(values_only=True):
            writer.writerow(["" if value is None else value for value in row])
    return file_result(output, "siaq-excel-to-csv.csv", "text/csv; charset=utf-8")


def text_to_word(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    document = Document()
    for line in paths[0].read_text(encoding="utf-8", errors="replace").splitlines():
        document.add_paragraph(line)
    output = workdir / "text.docx"
    document.save(output)
    return file_result(output, "siaq-text-to-word.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


def word_to_text(paths: list[Path], workdir: Path, _options: dict) -> ProcessResult:
    document = Document(_modern_word(paths[0], workdir))
    lines = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        lines.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
    output = workdir / "word.txt"
    output.write_text("\n".join(lines), encoding="utf-8")
    return file_result(output, "siaq-word-to-text.txt", "text/plain; charset=utf-8")


OFFICE_PROCESSORS = {
    "pdf-to-word": pdf_to_word,
    "word-to-pdf": office_to_pdf,
    "pdf-to-excel": pdf_to_excel,
    "excel-to-pdf": office_to_pdf,
    "pdf-to-powerpoint": pdf_to_powerpoint,
    "powerpoint-to-pdf": office_to_pdf,
    "excel-to-word": excel_to_word,
    "word-to-excel": word_to_excel,
    "word-to-powerpoint": word_to_powerpoint,
    "powerpoint-to-word": powerpoint_to_word,
    "excel-to-powerpoint": excel_to_powerpoint,
    "powerpoint-to-excel": powerpoint_to_excel,
    "csv-to-excel": csv_to_excel,
    "excel-to-csv": excel_to_csv,
    "text-to-word": text_to_word,
    "word-to-text": word_to_text,
}
OFFICE_PROCESSORS.update({slug: (lambda p, w, o, tool_slug=slug: office_format(p, w, o, tool_slug)) for slug in FORMAT_TARGETS})
